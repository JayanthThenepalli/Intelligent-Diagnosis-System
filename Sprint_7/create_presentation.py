import os
import subprocess
import sys

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_title_box(slide, text, top=Inches(0.5), height=Inches(0.8)):
    title_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9.0), height)
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)  # Deep Navy #1E3A8A
    return title_box

def add_bullets_box(slide, left, top, width, height, bullet_points):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    for i, pt in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pt
        p.font.name = 'Calibri'
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 65, 85)  # Slate Gray #334155
        p.space_after = Pt(8)
    return textbox

def main():
    prs = Presentation()
    
    # Configure slide dimensions to standard 16:9 widescreen layout
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)
    
    blank_layout = prs.slide_layouts[6]
    bg_color = RGBColor(248, 250, 252)  # Light Slate White #F8FAFC

    # ----------------------------------------------------
    # Slide 1: Title Slide
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    
    # Add title container
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "MEDIWISE AI"
    p.font.name = 'Calibri'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "An Intelligent Clinical Diagnosis Portal for Triage and Lesion Classification Under Uncertainty"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(59, 130, 246) # Primary Accent #3B82F6
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(14)
    
    p3 = tf.add_paragraph()
    p3.text = "Presenter: Jayanth Thenepalle | EBTM 881 Capstone Project"
    p3.font.name = 'Calibri'
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(100, 116, 139) # Secondary Accent #64748B
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "Advisor: Dr. Faculty Supervisor"
    p4.font.name = 'Calibri'
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(100, 116, 139)
    p4.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # Slide 2: Introduction & Motivation
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "1. Introduction & Motivation")
    bullets = [
        "Clinical Access Bottlenecks: Administrative delays frequently delay preliminary medical triage screenings, risking patient health.",
        "Uncertainty Constraints: Symptom assessments often deal with sparse or incomplete client input variables.",
        "System Decoupling: Serving models centrally on a secure Python server prevents IP exposure and client resource fatigue.",
        "Diagnostic Modalities: Implements both a tabular symptom checker and image-based dermoscopic cancer lesion classifications."
    ]
    add_bullets_box(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5), bullets)

    # ----------------------------------------------------
    # Slide 3: Research Objectives
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "2. Research Objectives")
    bullets = [
        "Stateless API Infrastructure: Design a horizontally scalable backend server capable of handling parallel client requests.",
        "High-Precision Classifier: Train and validate neural network architectures capable of triaging medical parameters.",
        "Cloud Audit Logging: Integrate automated database tracking to log session diagnostics statelessly.",
        "Production Quality Engineering: Implement modern CORS security middleware, Docker containerization, and automated CI/CD deployments."
    ]
    add_bullets_box(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5), bullets)

    # ----------------------------------------------------
    # Slide 4: Literature Review Summary Matrix
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "3. Literature Review Matrix")
    
    # Add table summary box
    table_shape = slide.shapes.add_table(5, 4, Inches(0.5), Inches(1.4), Inches(9.0), Inches(3.8))
    table = table_shape.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.5)
    
    headers = ['Author (Year)', 'Methodology', 'Key Findings', 'MediWise Connection']
    for idx, name in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 58, 138)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            
    lit_rows = [
        ["Tschandl (2018)", "HAM10000 dataset collection.", "10k pigmented lesion scans.", "Extended via Transfer Learning MobileNetV2 CNN."],
        ["Sandler (2018)", "MobileNetV2 inverted residuals.", "Optimizes edge parameter footprint.", "Reduces model footprint in cloud server to ~340MB."],
        ["Esteva (2017)", "Deep CNN lesion scanning.", "Deep ML matches professional clinicians.", "Deploys a stateless, web-accessible scanner portal."],
        ["Chen (2020)", "CDSS using relational databases.", "DB locks bottleneck concurrent runs.", "Implements asynchronous Motor client, removing blocks."]
    ]
    for row_idx, data in enumerate(lit_rows):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9.5)
                p.font.color.rgb = RGBColor(51, 65, 85)

    # ----------------------------------------------------
    # Slide 5: Clinical Workflow Pipeline
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "4. Clinical Workflow Pipeline")
    bullets = [
        "Interactive Triage Path: Patients input symptom lists or upload lesion images via a responsive client dashboard.",
        "FastAPI Server Validation: Validates incoming request payloads using strict Pydantic schemas.",
        "Model Triage Execution: Dispatches payloads to local ANN or CNN instances cached in server memory.",
        "Asynchronous Cloud Logging: Motor driver schedules MongoDB insertions without blocking the client response thread."
    ]
    add_bullets_box(slide, Inches(0.5), Inches(1.4), Inches(4.5), Inches(3.8), bullets)
    
    # Insert Workflow flowchart
    workflow_img = "Clinical_Workflow.png"
    if os.path.exists(workflow_img):
        slide.shapes.add_picture(workflow_img, Inches(5.2), Inches(1.4), width=Inches(4.3))

    # ----------------------------------------------------
    # Slide 6: System Architecture
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "5. Three-Tier System Architecture")
    bullets = [
        "Presentation (Vercel): React.js SPA utilizing custom CSS tokens, EKG waves, and radial confidence dials.",
        "Service (Render): Stateless FastAPI ASGI web app serving model instances cached dynamically in RAM.",
        "Persistence (MongoDB Atlas): Schema-less cloud replica set collection logging diagnostic session data asynchronously."
    ]
    add_bullets_box(slide, Inches(0.5), Inches(1.4), Inches(4.5), Inches(3.8), bullets)
    
    # Insert System Architecture diagram
    arch_img = "System_Architecture.png"
    if os.path.exists(arch_img):
        slide.shapes.add_picture(arch_img, Inches(5.2), Inches(1.4), width=Inches(4.3))

    # ----------------------------------------------------
    # Slide 7: System Integration Workflow (DFD)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "6. System Integration Workflow (DFD)")
    bullets = [
        "HTTPS Client Dispatch: React form triggers Axios HTTP POST requests conveying symptom vectors or image FormData.",
        "FastAPI Router Middleware: Enforces CORS constraints, validates models, and processes predictions.",
        "Model Memory Load: Caches network weights during startup to run calculations entirely in RAM.",
        "Non-blocking Persistence: Motor schedules insert_one() tasks concurrently before sending JSON predictions back."
    ]
    add_bullets_box(slide, Inches(0.5), Inches(1.4), Inches(4.5), Inches(3.8), bullets)
    
    # Insert Integration DFD
    dfd_img = "System_Integration_Workflow.png"
    if os.path.exists(dfd_img):
        slide.shapes.add_picture(dfd_img, Inches(5.2), Inches(1.4), width=Inches(4.3))

    # ----------------------------------------------------
    # Slide 8: Model Engineering & Comparative Matrix
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "7. Comparative Model Performance Matrix")
    
    # Add comparison table
    table_shape = slide.shapes.add_table(5, 6, Inches(0.5), Inches(1.4), Inches(9.0), Inches(3.5))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(1.4)
    table.columns[2].width = Inches(1.4)
    table.columns[3].width = Inches(1.3)
    table.columns[4].width = Inches(1.3)
    table.columns[5].width = Inches(1.4)
    
    cols = ['Model Architecture', 'Target', 'Accuracy', 'Precision', 'Recall', 'F1-Score']
    for idx, name in enumerate(cols):
        cell = table.cell(0, idx)
        cell.text = name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 58, 138)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            
    matrix_rows = [
        ['Decision Tree (Baseline)', 'Symptoms', '92.30%', '92.45%', '92.30%', '92.31%'],
        ['Random Forest (Baseline)', 'Symptoms', '96.50%', '96.60%', '96.50%', '96.52%'],
        ['Artificial Neural Network', 'Symptoms', '98.42%', '98.44%', '98.42%', '98.42%'],
        ['MobileNetV2 CNN', 'Lesion Scans', '89.50%', '88.90%', '89.50%', '88.70%']
    ]
    for row_idx, data in enumerate(matrix_rows):
        for col_idx, text in enumerate(data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(51, 65, 85)

    # ----------------------------------------------------
    # Slide 9: Baseline Classifiers Performance
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "8. Baseline Classifier Evaluation")
    
    dt_img = "Decision_Tree_Confusion_Matrix.png"
    rf_img = "Random_Forest_Confusion_Matrix.png"
    
    if os.path.exists(dt_img):
        slide.shapes.add_picture(dt_img, Inches(0.5), Inches(1.4), width=Inches(4.2))
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(4.2), Inches(0.5))
        tbox.text_frame.text = "Decision Tree Baseline (Accuracy: 92.30%)"
        tbox.text_frame.paragraphs[0].font.size = Pt(11)
        tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    if os.path.exists(rf_img):
        slide.shapes.add_picture(rf_img, Inches(5.3), Inches(1.4), width=Inches(4.2))
        tbox = slide.shapes.add_textbox(Inches(5.3), Inches(4.8), Inches(4.2), Inches(0.5))
        tbox.text_frame.text = "Random Forest Ensemble Baseline (Accuracy: 96.50%)"
        tbox.text_frame.paragraphs[0].font.size = Pt(11)
        tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # Slide 10: Proposed Symptom ANN Performance
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "9. Proposed Symptom Predictor ANN")
    
    ann_img = "ANN_Confusion_Matrix.png"
    ann_hist = "ANN_Training_History.png"
    
    if os.path.exists(ann_img):
        slide.shapes.add_picture(ann_img, Inches(0.5), Inches(1.4), width=Inches(4.2))
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(4.2), Inches(0.5))
        tbox.text_frame.text = "ANN Confusion Matrix (Accuracy: 98.42%)"
        tbox.text_frame.paragraphs[0].font.size = Pt(11)
        tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    if os.path.exists(ann_hist):
        slide.shapes.add_picture(ann_hist, Inches(5.3), Inches(1.4), width=Inches(4.2))
        tbox = slide.shapes.add_textbox(Inches(5.3), Inches(4.8), Inches(4.2), Inches(0.5))
        tbox.text_frame.text = "ANN Loss & Accuracy Training History Curves"
        tbox.text_frame.paragraphs[0].font.size = Pt(11)
        tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # Slide 11: Proposed Skin Lesion CNN Performance
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "10. Proposed Skin Lesion CNN")
    
    cnn_img = "CNN_Confusion_Matrix.png"
    cnn_hist = "CNN_Training_History.png"
    
    if os.path.exists(cnn_img):
        slide.shapes.add_picture(cnn_img, Inches(0.5), Inches(1.4), width=Inches(4.2))
        tbox = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(4.2), Inches(0.5))
        tbox.text_frame.text = "CNN Confusion Matrix (Accuracy: 89.50%)"
        tbox.text_frame.paragraphs[0].font.size = Pt(11)
        tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    if os.path.exists(cnn_hist):
        slide.shapes.add_picture(cnn_hist, Inches(5.3), Inches(1.4), width=Inches(4.2))
        tbox = slide.shapes.add_textbox(Inches(5.3), Inches(4.8), Inches(4.2), Inches(0.5))
        tbox.text_frame.text = "CNN Loss & Accuracy Training Curves"
        tbox.text_frame.paragraphs[0].font.size = Pt(11)
        tbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # Slide 12: Software Engineering Innovations
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "11. Software Engineering Innovations")
    bullets = [
        "Specialist Routing logic: Custom algorithms parsing diagnostic classifications to assign patients to appropriate specialists.",
        "Print-Media CSS referral exporter: Hides navigation panels, generating clean clinical referral letters.",
        "Asynchronous Database Client: Leverages the motor driver to coordinate connection queues asynchronously.",
        "Dockerized footprint optimization: Isolates tensorflow-cpu in multi-stage builds, reducing image size from >3GB to ~450MB."
    ]
    add_bullets_box(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5), bullets)

    # ----------------------------------------------------
    # Slide 13: CI/CD Cloud Deployments & Verification
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "12. Production Deployments & Health")
    bullets = [
        "Vite Frontend Client (Vercel): Global Edge Network serving React code with SSL active (load times < 200ms).",
        "FastAPI Backend API (Render): Stateless docker container online with memory consumption stabilized at ~340MB.",
        "MongoDB Atlas cloud: Active connection replica set writing diagnostic session details with timestamps.",
        "Low Inference Latency: Average response delays measured at 85ms for ANN symptom check and 120ms for CNN scans."
    ]
    add_bullets_box(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5), bullets)

    # ----------------------------------------------------
    # Slide 14: Conclusion & Future Work
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide, bg_color)
    create_title_box(slide, "13. Conclusion & Future Scope")
    bullets = [
        "Conclusion: Successful design, integration, and cloud deployment of an intelligent diagnostic suite resolving clinical uncertainty.",
        "Future Work - Horizontal Scaling: Introduce Celery distributed worker tasks for deep neural scans under heavier traffic load.",
        "Future Work - Active Integrations: Bind model outputs directly into FHIR clinical database structures.",
        "Future Work - Explainable AI: Integrate Grad-CAM heatmaps showing clinicians exactly which image regions prompted predictions."
    ]
    add_bullets_box(slide, Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5), bullets)

    output_path = "EBTM881_Capstone_Final_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully as {output_path}!")

if __name__ == '__main__':
    main()
